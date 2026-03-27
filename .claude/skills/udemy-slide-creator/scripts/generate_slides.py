#!/usr/bin/env python3
"""
Generate Dyer Innovation branded PowerPoint slides for Udemy courses.

Takes parsed lecture JSON (from parse_lecture.py) and produces a .pptx file
using the Dyer Innovation design system.

Usage:
    python generate_slides.py <parsed_json_path> <logo_path> <output_pptx_path>

Example:
    python generate_slides.py /tmp/lecture_parsed.json \
        /Users/jonathandyer/Documents/Dyer_Innovation/dev/utils/dyer_innovation_logo.png \
        course/slides/lecture-2.3-agentic-loops.pptx
"""

import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("python-pptx not found. Install with: pip3 install python-pptx")
    sys.exit(1)


# ---------------------------------------------------------------------------
# Brand color constants
# ---------------------------------------------------------------------------
C_LEAF_GREEN   = RGBColor(0x3C, 0xAF, 0x50)  # Primary green
C_DEEP_GREEN   = RGBColor(0x1B, 0x8A, 0x5A)  # Gradient / accent
C_TEAL         = RGBColor(0x0D, 0x73, 0x77)  # Teal accent
C_DARK_NAVY    = RGBColor(0x1A, 0x3A, 0x4A)  # Headings, dark bg
C_LIGHT_MINT   = RGBColor(0xE8, 0xF5, 0xEB)  # Content slide bg
C_AMBER        = RGBColor(0xE3, 0xA0, 0x08)  # Exam tips
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_NEAR_BLACK   = RGBColor(0x11, 0x19, 0x28)  # Body text
C_CODE_DARK    = RGBColor(0x1E, 0x2D, 0x3D)  # Code bg
C_CODE_LIGHT   = RGBColor(0xA8, 0xD5, 0xC2)  # Code text
C_CODE_DARKER  = RGBColor(0x15, 0x25, 0x35)  # Code inner box bg
C_TRAP_BG      = RGBColor(0xFF, 0xF0, 0xF0)  # Exam trap box
C_CORRECT_BG   = RGBColor(0xF0, 0xFF, 0xF4)  # Exam correct box
C_TRAP_LABEL   = RGBColor(0xC0, 0x39, 0x2B)  # Trap label red
C_CORRECT_LABEL = RGBColor(0x1B, 0x8A, 0x5A) # Correct label green

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Logo placement
LOGO_CORNER_W = Inches(0.8)
LOGO_CORNER_X = Inches(12.333)   # 13.333 - 0.8 - 0.2
LOGO_CORNER_Y = Inches(6.85)     # 7.5 - 0.5 - 0.15

LOGO_TITLE_W  = Inches(1.5)
LOGO_TITLE_X  = Inches(5.917)    # centered: (13.333 - 1.5) / 2
LOGO_TITLE_Y  = Inches(6.75)

# Title bar height
TITLE_BAR_H = Inches(1.1)


# ---------------------------------------------------------------------------
# Helper utilities
# ---------------------------------------------------------------------------

def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    """Add a blank slide (layout 6)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color: RGBColor):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor, no_line=True):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if no_line:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_name="Calibri", font_size=20,
                bold=False, color: RGBColor = C_NEAR_BLACK,
                align=PP_ALIGN.LEFT, word_wrap=True, italic=False) -> object:
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullets(slide, x, y, w, h, bullets: list, font_size=20,
                font_name="Calibri", color: RGBColor = C_NEAR_BLACK,
                bullet_color: RGBColor = None):
    """Add a text frame with multiple bullet paragraphs."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.space_before = Pt(4)
        run = p.runs[0] if p.runs else p.add_run()
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_logo(slide, x, y, w, logo_path: str):
    if logo_path and Path(logo_path).exists():
        slide.shapes.add_picture(str(logo_path), x, y, width=w)


def set_notes(slide, text: str):
    if text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text


def extract_bullets_from_narration(narration: str) -> list:
    """Extract bullet items from narration text, or split into sentences."""
    lines = [l.strip() for l in narration.splitlines() if l.strip()]
    bullets = []
    for line in lines:
        clean = re.sub(r'^[-*•]\s*', '', line)
        if clean and not clean.startswith('['):
            bullets.append(clean)
    # Limit to 6 bullets for readability
    return bullets[:6] if bullets else [narration[:200]]


def narration_to_bullets(segments: list) -> list:
    """Convert slide segments into bullet strings for the slide body."""
    all_bullets = []
    for seg in segments:
        text = seg.get('narration', '')
        all_bullets.extend(extract_bullets_from_narration(text))
    return all_bullets[:6]


# ---------------------------------------------------------------------------
# Slide layout builders
# ---------------------------------------------------------------------------

def build_title_slide(prs, slide_data: dict, logo_path: str):
    slide = blank_slide(prs)
    set_bg(slide, C_DARK_NAVY)

    # Green accent bar at top
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), C_LEAF_GREEN)

    # Main title — centered vertically
    title = slide_data['title']
    add_textbox(slide,
                Inches(1), Inches(2.2), Inches(11.333), Inches(2),
                title, font_size=44, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

    # Subtitle from metadata (lecture number + section)
    meta = slide_data.get('_meta', {})
    subtitle_parts = []
    if meta.get('lecture_number'):
        subtitle_parts.append(f"Lecture {meta['lecture_number']}")
    if meta.get('section'):
        subtitle_parts.append(meta['section'])
    subtitle = '  •  '.join(subtitle_parts)
    if subtitle:
        add_textbox(slide,
                    Inches(1), Inches(4.4), Inches(11.333), Inches(0.6),
                    subtitle, font_size=20, color=RGBColor(0xB0, 0xC8, 0xD0),
                    align=PP_ALIGN.CENTER)

    # Logo centered bottom
    add_logo(slide, LOGO_TITLE_X, LOGO_TITLE_Y, LOGO_TITLE_W, logo_path)
    set_notes(slide, slide_data.get('speaker_notes', ''))
    return slide


def build_content_slide(prs, slide_data: dict, logo_path: str):
    slide = blank_slide(prs)
    set_bg(slide, C_LIGHT_MINT)

    # Left green accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, C_LEAF_GREEN)

    # Title bar
    add_rect(slide, Inches(0.08), Inches(0), Inches(13.253), TITLE_BAR_H, C_DARK_NAVY)
    add_textbox(slide,
                Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.9),
                slide_data['title'], font_size=32, bold=True, color=C_WHITE)

    # Body bullets
    bullets = narration_to_bullets(slide_data.get('segments', []))
    if bullets:
        add_bullets(slide,
                    Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5),
                    bullets, font_size=20, color=C_NEAR_BLACK)

    add_logo(slide, LOGO_CORNER_X, LOGO_CORNER_Y, LOGO_CORNER_W, logo_path)
    set_notes(slide, slide_data.get('speaker_notes', ''))
    return slide


def build_code_slide(prs, slide_data: dict, logo_path: str):
    slide = blank_slide(prs)
    set_bg(slide, C_CODE_DARK)

    # Teal accent line below title area
    add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), C_TEAL)

    # Title
    add_textbox(slide,
                Inches(0.5), Inches(0.1), Inches(12.3), Inches(0.95),
                slide_data['title'], font_size=32, bold=True, color=C_WHITE)

    # Code block background box
    code_text = slide_data.get('code_block') or ''
    if not code_text:
        # Fall back to narration
        narration = ' '.join(s.get('narration', '') for s in slide_data.get('segments', []))
        code_match = re.search(r'```[\w]*\n(.*?)```', narration, re.DOTALL)
        code_text = code_match.group(1).strip() if code_match else narration[:400]

    add_rect(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(5.7), C_CODE_DARKER)

    # Code text
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = code_text
    run = p.runs[0] if p.runs else p.add_run()
    run.font.name = "Courier New"
    run.font.size = Pt(14)
    run.font.color.rgb = C_CODE_LIGHT

    add_logo(slide, LOGO_CORNER_X, LOGO_CORNER_Y, LOGO_CORNER_W, logo_path)
    set_notes(slide, slide_data.get('speaker_notes', ''))
    return slide


def build_comparison_slide(prs, slide_data: dict, logo_path: str):
    slide = blank_slide(prs)
    set_bg(slide, C_WHITE)

    # Title bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, TITLE_BAR_H, C_DARK_NAVY)
    add_textbox(slide,
                Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.9),
                slide_data['title'], font_size=32, bold=True, color=C_WHITE)

    # Vertical divider
    add_rect(slide, Inches(6.617), Inches(1.2), Inches(0.05), Inches(5.8), C_LEAF_GREEN)

    # Left column
    segments = slide_data.get('segments', [])
    left_text = segments[0]['narration'] if segments else ''
    right_text = segments[1]['narration'] if len(segments) > 1 else ''

    add_textbox(slide, Inches(0.3), Inches(1.2), Inches(5.9), Inches(0.5),
                "Option A", font_size=18, bold=True, color=C_LEAF_GREEN)
    left_bullets = extract_bullets_from_narration(left_text)
    add_bullets(slide, Inches(0.3), Inches(1.8), Inches(5.9), Inches(5.0),
                left_bullets, font_size=18, color=C_NEAR_BLACK)

    # Right column
    add_textbox(slide, Inches(6.9), Inches(1.2), Inches(6.0), Inches(0.5),
                "Option B", font_size=18, bold=True, color=C_TEAL)
    right_bullets = extract_bullets_from_narration(right_text)
    add_bullets(slide, Inches(6.9), Inches(1.8), Inches(6.0), Inches(5.0),
                right_bullets, font_size=18, color=C_NEAR_BLACK)

    add_logo(slide, LOGO_CORNER_X, LOGO_CORNER_Y, LOGO_CORNER_W, logo_path)
    set_notes(slide, slide_data.get('speaker_notes', ''))
    return slide


def build_exam_tip_slide(prs, slide_data: dict, logo_path: str):
    slide = blank_slide(prs)
    set_bg(slide, C_WHITE)

    # Amber banner at top
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.3), C_AMBER)
    add_textbox(slide,
                Inches(0.5), Inches(0.15), Inches(12.333), Inches(1.0),
                "EXAM TIP", font_size=36, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)

    # Title / exam tip summary
    add_textbox(slide,
                Inches(0.5), Inches(1.45), Inches(12.333), Inches(0.7),
                slide_data['title'], font_size=22, bold=True,
                color=C_DARK_NAVY, align=PP_ALIGN.CENTER)

    narration = slide_data.get('speaker_notes', '')
    # Try to extract Trap vs Correct from narration
    trap_match = re.search(r'[Tt]rap[:\s]+(.+?)(?:[Cc]orrect|$)', narration, re.DOTALL)
    correct_match = re.search(r'[Cc]orrect\s+[Aa]pproach[:\s]+(.+?)(?:\n\n|$)', narration, re.DOTALL)
    trap_text = trap_match.group(1).strip()[:300] if trap_match else narration[:200]
    correct_text = correct_match.group(1).strip()[:300] if correct_match else ""

    # Trap box (left)
    add_rect(slide, Inches(0.4), Inches(2.3), Inches(5.8), Inches(4.5), C_TRAP_BG)
    add_textbox(slide, Inches(0.6), Inches(2.4), Inches(5.4), Inches(0.5),
                "The Trap", font_size=18, bold=True, color=C_TRAP_LABEL)
    add_textbox(slide, Inches(0.6), Inches(3.0), Inches(5.4), Inches(3.5),
                trap_text, font_size=16, color=C_NEAR_BLACK, word_wrap=True)

    # Correct box (right)
    add_rect(slide, Inches(7.133), Inches(2.3), Inches(5.8), Inches(4.5), C_CORRECT_BG)
    add_textbox(slide, Inches(7.333), Inches(2.4), Inches(5.4), Inches(0.5),
                "Correct Approach", font_size=18, bold=True, color=C_CORRECT_LABEL)
    add_textbox(slide, Inches(7.333), Inches(3.0), Inches(5.4), Inches(3.5),
                correct_text, font_size=16, color=C_NEAR_BLACK, word_wrap=True)

    add_logo(slide, LOGO_CORNER_X, LOGO_CORNER_Y, LOGO_CORNER_W, logo_path)
    set_notes(slide, narration)
    return slide


def build_takeaway_slide(prs, slide_data: dict, logo_path: str):
    slide = blank_slide(prs)

    # Gradient-effect: two overlapping rects (pptx doesn't support true gradients simply)
    set_bg(slide, C_DARK_NAVY)
    add_rect(slide, Inches(0), Inches(5.0), SLIDE_W, Inches(2.5), C_TEAL)
    # Overlay top to blend
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(5.2),
             RGBColor(0x1A, 0x3A, 0x4A))

    # Title
    add_textbox(slide,
                Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.9),
                "What to Remember", font_size=36, bold=True,
                color=C_WHITE, align=PP_ALIGN.LEFT)

    # Green separator line
    add_rect(slide, Inches(0.5), Inches(1.15), Inches(12.333), Inches(0.04), C_LEAF_GREEN)

    # Bullet points
    bullets = narration_to_bullets(slide_data.get('segments', []))
    y = Inches(1.35)
    for bullet in bullets:
        # Green dot marker
        add_rect(slide, Inches(0.5), y + Inches(0.1), Inches(0.12), Inches(0.12), C_LEAF_GREEN)
        add_textbox(slide, Inches(0.75), y, Inches(12.0), Inches(0.65),
                    bullet, font_size=20, color=C_WHITE)
        y += Inches(0.72)

    add_logo(slide, LOGO_CORNER_X, LOGO_CORNER_Y, LOGO_CORNER_W, logo_path)
    set_notes(slide, slide_data.get('speaker_notes', ''))
    return slide


def build_section_slide(prs, slide_data: dict, logo_path: str):
    slide = blank_slide(prs)

    # Green gradient background (two rects)
    set_bg(slide, C_LEAF_GREEN)
    add_rect(slide, Inches(0), Inches(4.0), SLIDE_W, Inches(3.5), C_DEEP_GREEN)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(4.2), C_LEAF_GREEN)

    # Section title
    add_textbox(slide,
                Inches(1), Inches(2.5), Inches(11.333), Inches(2.0),
                slide_data['title'], font_size=40, bold=True,
                color=C_WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    narration = ' '.join(s.get('narration', '') for s in slide_data.get('segments', []))
    subtitle = narration[:120] if narration else ''
    if subtitle:
        add_textbox(slide,
                    Inches(1.5), Inches(4.7), Inches(10.333), Inches(0.7),
                    subtitle, font_size=18, italic=True,
                    color=RGBColor(0xD0, 0xF0, 0xD8), align=PP_ALIGN.CENTER)

    add_logo(slide, LOGO_TITLE_X, LOGO_TITLE_Y, Inches(1.2), logo_path)
    set_notes(slide, slide_data.get('speaker_notes', ''))
    return slide


def build_quiz_slide(prs, slide_data: dict, logo_path: str):
    slide = blank_slide(prs)
    set_bg(slide, C_LIGHT_MINT)

    # Teal title bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, TITLE_BAR_H, C_TEAL)
    add_textbox(slide,
                Inches(0.3), Inches(0.1), Inches(12.7), Inches(0.9),
                "Knowledge Check", font_size=28, bold=True, color=C_WHITE)

    # Question text
    question = slide_data['title']
    add_textbox(slide,
                Inches(0.5), Inches(1.3), Inches(12.333), Inches(1.2),
                question, font_size=22, bold=True, color=C_DARK_NAVY,
                word_wrap=True)

    # Answer options as boxes
    bullets = narration_to_bullets(slide_data.get('segments', []))
    labels = ['A', 'B', 'C', 'D']
    positions = [
        (Inches(0.5), Inches(2.8)),
        (Inches(6.9), Inches(2.8)),
        (Inches(0.5), Inches(4.5)),
        (Inches(6.9), Inches(4.5)),
    ]
    for i, (bx, by) in enumerate(positions):
        if i >= len(bullets):
            break
        add_rect(slide, bx, by, Inches(6.0), Inches(1.4), C_WHITE)
        add_textbox(slide, bx + Inches(0.15), by + Inches(0.1), Inches(0.5), Inches(1.1),
                    labels[i], font_size=20, bold=True, color=C_TEAL)
        add_textbox(slide, bx + Inches(0.7), by + Inches(0.1), Inches(5.1), Inches(1.1),
                    bullets[i], font_size=16, color=C_NEAR_BLACK, word_wrap=True)

    add_logo(slide, LOGO_CORNER_X, LOGO_CORNER_Y, LOGO_CORNER_W, logo_path)
    set_notes(slide, slide_data.get('speaker_notes', ''))
    return slide


# ---------------------------------------------------------------------------
# [click] segment splitting
# ---------------------------------------------------------------------------

def expand_click_slides(slides: list) -> list:
    """
    Expand slides with [click] markers into multiple slides.
    Each segment after a [click] becomes a new slide with the same type and title.
    """
    expanded = []
    for slide_data in slides:
        if not slide_data.get('has_clicks'):
            expanded.append(slide_data)
            continue

        segments = slide_data.get('segments', [])
        if len(segments) <= 1:
            expanded.append(slide_data)
            continue

        # Create one slide per segment
        base_title = slide_data['title']
        for i, seg in enumerate(segments):
            new_slide = dict(slide_data)
            new_slide['segments'] = [seg]
            new_slide['speaker_notes'] = seg['narration']
            if i > 0:
                new_slide['title'] = f"{base_title} (cont.)"
            new_slide['has_clicks'] = False
            expanded.append(new_slide)

    return expanded


# ---------------------------------------------------------------------------
# Main dispatch
# ---------------------------------------------------------------------------

LAYOUT_BUILDERS = {
    'TITLE':      build_title_slide,
    'CONTENT':    build_content_slide,
    'CODE':       build_code_slide,
    'COMPARISON': build_comparison_slide,
    'EXAM_TIP':   build_exam_tip_slide,
    'TAKEAWAY':   build_takeaway_slide,
    'SECTION':    build_section_slide,
    'QUIZ':       build_quiz_slide,
}


def generate_presentation(parsed: dict, logo_path: str, output_path: str):
    prs = new_presentation()
    metadata = parsed.get('metadata', {})
    slides = expand_click_slides(parsed.get('slides', []))

    print(f"Generating {len(slides)} slides...")

    for slide_data in slides:
        # Attach metadata for use by title slide builder
        slide_data['_meta'] = metadata
        slide_type = slide_data.get('slide_type', 'CONTENT')
        builder = LAYOUT_BUILDERS.get(slide_type, build_content_slide)
        builder(prs, slide_data, logo_path)
        print(f"  Slide {slide_data['slide_number']}: {slide_data['title']} [{slide_type}]")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"\nSaved: {output}")
    return str(output)


def main():
    if len(sys.argv) < 4:
        print("Usage: generate_slides.py <parsed_json> <logo_path> <output_pptx>")
        sys.exit(1)

    json_path, logo_path, output_path = sys.argv[1], sys.argv[2], sys.argv[3]

    with open(json_path) as f:
        parsed = json.load(f)

    generate_presentation(parsed, logo_path, output_path)


if __name__ == '__main__':
    main()
